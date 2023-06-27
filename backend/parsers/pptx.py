import tempfile
from fastapi import UploadFile
from pptx import Presentation


async def pptx(file: UploadFile) -> str:
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        temp_file.write(await file.read())
        temp_file.seek(0)

        prs = Presentation(temp_file.name)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content)
